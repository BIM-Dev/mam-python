from pptx import Presentation
from PIL import Image
from pptx.util import Cm
from pptx.chart.data import CategoryChartData
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import pandas
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from datetime import datetime
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_ANCHOR
import psycopg2
from psycopg2.extras import DictCursor
import requests
from io import BytesIO
from decimal import Decimal, ROUND_HALF_UP
from PIL import Image
from django.conf import settings
connection_parameters = {
          # "host": "10.4.47.220",
          #   "port": "5432",
          #   "database": "DataAnalysis",
          #   "user": "postgres",
          #   "password": "aa4f8ef88f3b1e65aee010a79e69d71b"
"user" : "postgres",
"password" : "docker",
"host": settings.DBHOST,
"port" : "5432",
"database": "generic_am_mam"
    }
cn_template_dir = r'.\mam_api_app\service\Impression_Report_TemplateCN.pptx'
en_template_dir = r'.\mam_api_app\service\Impression_Report_TemplateEN.pptx'
schema = 'public'
class generateReport:
    # def __init__(self):
    #     self.connection_parameters = {
    #         "host": "10.4.47.220",
    #         "port": "5432",
    #         "database": "DataAnalysis",
    #         "user": "postgres",
    #         "password": "aa4f8ef88f3b1e65aee010a79e69d71b"
    #     }
    #     self.schema = 'mam_impression_report_test'
    #     self.cn_template_dir = r'\mam_api_app\service\Impression_Report_TemplateCN.pptx'
    #     self.en_template_dir = r'C:\Users\zewen.liang\OneDrive - JCDECAUX\桌面\Impression Report Template - EN.pptx'
    #

    def __init__(self):
        pass

    def get_report_info(self, report_id, env_info_code):
        conn = psycopg2.connect(**connection_parameters)
        cur = conn.cursor(cursor_factory=DictCursor)
        report_info_sql = f"""
        SELECT no, env_info_code, total_impression, campaign_name_cn, campaign_name_en, 
                campaign_start_date, campaign_end_date, client_name_cn, client_name_en, client_logo_cn, 
                client_logo_en, total_cost, cpm, env_company_id
	        FROM {schema}.impression_report
	    WHERE is_active = TRUE and id = {report_id} and env_info_code = '{env_info_code}'
        """
        cur.execute(report_info_sql)
        report_info_data = cur.fetchall()
        print(report_info_data)
        if len(report_info_data) != 1:
            raise ValueError("The number of report_info_data rows is not equal to 1.")
        else:
            dict_report_info_data = dict(report_info_data[0])

        asset_provider_sql = f"""
        SELECT id, env_info_code, name_cn, name_en, logo, full_name_cn, full_name_en
	        FROM {schema}.env_company
	    WHERE is_active = TRUE and id = {dict_report_info_data['env_company_id']} and env_info_code = '{env_info_code}'
        """
        cur.execute(asset_provider_sql)
        asset_provider_data = cur.fetchall()
        if len(asset_provider_data) != 1:
            raise ValueError("The number of asset_provider_data rows is not equal to 1.")
        else:
            dict_asset_provider_data = dict(asset_provider_data[0])

        report_item_sql = f"""
        SELECT date, sot, line_code, line_name_cn, line_name_en, station_code, station_name_cn, station_name_en, asset_code, asset_name_cn, asset_name_en, impression
	        FROM public.impression_report_item
	    WHERE env_info_code = '{env_info_code}' and impression_report_id = {report_id}
        """
        cur.execute(report_item_sql)
        report_item_data = cur.fetchall()
        df_report_item_data = pandas.DataFrame(report_item_data, columns=[desc[0] for desc in cur.description])

        cur.close()
        conn.close()
        return dict_report_info_data, dict_asset_provider_data, df_report_item_data

    def insert_image_into_placeholder(self, placeholder, image_dir, type, vertical_align="top", horizontal_align="left"):
        # image = Image.open(BytesIO(requests.get(image_dir).content))
        if image_dir is None:
            return
        # response = requests.get(image_dir)
        # image = Image.open(BytesIO(response.content))
        if type == 'clientLogo':
            image = Image.open(image_dir)
        else:
            response = requests.get(image_dir)
            image = Image.open(BytesIO(response.content))
        # print(placeholder)
        # image = Image.open(image_dir)


        image_width, image_height = image.size
        placeholder_width = placeholder.width.cm
        placeholder_height = placeholder.height.cm
        placeholder_left = placeholder.left.cm
        placeholder_top = placeholder.top.cm
        placeholder.width = Cm(placeholder_height * image_width / image_height)
        placeholder.height = Cm(placeholder_height)
        if horizontal_align == "left":
            left_deviation = 0
        elif horizontal_align == "right":
            left_deviation = placeholder.width.cm - placeholder_width
        if vertical_align == "top":
            top_deviation = 0
        elif vertical_align == "bottom":
            top_deviation = placeholder.height.cm - placeholder_height
        # pic = placeholder.insert_picture(image_dir)
        if type == 'companyLogo':
            pic = placeholder.insert_picture(BytesIO(response.content))
        else:
            pic =  placeholder.insert_picture(image_dir)

        pic.left = Cm(placeholder_left - left_deviation)
        pic.top = Cm(placeholder_top - top_deviation)
        pic.width = Cm(placeholder_height * image_width / image_height)
        pic.height = Cm(placeholder_height)

    def add_cover(self, presentation, layout, report_main_title, report_subtitle, campaign_period,
                  report_no, company_logo_dir, client_company_logo_dir):
        cover = layout
        prs = presentation
        prs.slides.add_slide(cover)
        slide = prs.slides[-1]
        placeholders = slide.placeholders
        placeholders[0].text = report_main_title
        placeholders[1].text = report_subtitle
        placeholders[14].text = campaign_period
        placeholders[18].text = report_no
        self.insert_image_into_placeholder(placeholders[16], company_logo_dir, type = 'companyLogo', horizontal_align="left", vertical_align="top")

        self.insert_image_into_placeholder(placeholders[17], client_company_logo_dir, type = 'clientLogo', horizontal_align="right", vertical_align="top")
        return prs

    def add_info_page(self, presentation, layout, title_of_this_page, descriptive_text, start_date, end_date, day_count, line_count, station_count,
                      asset_count, impression, cpm, company_logo_dir, client_company_logo_dir, report_no):
        prs = presentation
        prs.slides.add_slide(layout)
        slide = prs.slides[-1]
        placeholders = slide.placeholders
        placeholders[22].text = title_of_this_page
        placeholders[23].text = descriptive_text
        placeholders[26].text = start_date
        placeholders[28].text = end_date
        placeholders[30].text = day_count
        placeholders[32].text = line_count
        placeholders[34].text = station_count
        placeholders[36].text = asset_count
        placeholders[38].text = format(int(impression), ',')
        placeholders[38].text_frame.paragraphs[0].font.size = Pt(28) if int(impression) < 999999999 else Pt(24)
        placeholders[40].text = cpm
        self.insert_image_into_placeholder(placeholders[16], company_logo_dir, type='companyLogo', horizontal_align="left", vertical_align="top")
        self.insert_image_into_placeholder(placeholders[17], client_company_logo_dir, type='clientLogo', horizontal_align="right",
                                      vertical_align="top")
        placeholders[18].text = report_no
        return prs

    def add_confidentiality_agreement(self, company_logo_dir, client_company_logo_dir, report_no):
        prs.slides.add_slide(layouts[layout_index["Confidentiality Agreement"]])
        slide = prs.slides[-1]
        placeholders = slide.placeholders
        insert_image_into_placeholder(placeholders[16], company_logo_dir, type='companyLogo',  horizontal_align="left", vertical_align="top")
        insert_image_into_placeholder(placeholders[17], client_company_logo_dir, type='clientLogo',  horizontal_align="right",
                                      vertical_align="top")
        placeholders[18].text = report_no
        return prs

    def add_last_page(self, presentation, layout, company_logo_dir, client_company_logo_dir, report_no):
        prs = presentation
        prs.slides.add_slide(layout)
        slide = prs.slides[-1]
        placeholders = slide.placeholders
        self.insert_image_into_placeholder(placeholders[16], company_logo_dir, type='companyLogo', horizontal_align="left", vertical_align="top")
        self.insert_image_into_placeholder(placeholders[17], client_company_logo_dir, type='clientLogo', horizontal_align="right",
                                      vertical_align="top")
        placeholders[18].text = report_no
        return prs

    def add_chart_page(self, presentation, layout, title_of_this_page, descriptive_text, chart_or_table_1, chart_or_table_2, company_logo_dir, client_company_logo_dir, report_no):
        prs = presentation
        prs.slides.add_slide(layout)
        slide = prs.slides[-1]
        placeholders = slide.placeholders
        shapes = slide.shapes
        placeholders[22].text = title_of_this_page
        placeholders[23].text = descriptive_text
        if chart_or_table_1:
            if chart_or_table_1['object'] == 'chart':
                self.insert_chart_into_content_placeholder(placeholders[24], chart_or_table_1, shapes)
            elif chart_or_table_1['object'] == 'table':
                self.insert_table_into_content_placeholder(placeholders[24], chart_or_table_1, shapes)
        if chart_or_table_2:
            if chart_or_table_2['object'] == 'chart':
                self.insert_chart_into_content_placeholder(placeholders[25], chart_or_table_2, shapes)
            elif chart_or_table_2['object'] == 'table':
                self.insert_table_into_content_placeholder(placeholders[25], chart_or_table_2, shapes)

        self.insert_image_into_placeholder(placeholders[16], company_logo_dir, type='companyLogo',  horizontal_align="left", vertical_align="top")
        self.insert_image_into_placeholder(placeholders[17], client_company_logo_dir, type='clientLogo', horizontal_align="right",
                                      vertical_align="top")
        placeholders[18].text = report_no
        return prs

    def insert_chart_into_content_placeholder(self, placeholder, chart, shapes):
        idx = placeholder.placeholder_format.idx
        left = placeholder.left.cm
        top = placeholder.top.cm
        width = placeholder.width.cm
        height = placeholder.height.cm
        sp = placeholder.element
        sp.getparent().remove(sp)
        x, y, cx, cy = Cm(left), Cm(top), Cm(width), Cm(height)
        inserted_chart = shapes.add_chart(chart["type"], x, y, cx, cy, chart["data"]).chart
        if chart["type"] == XL_CHART_TYPE.LINE:
            inserted_chart.has_legend = False
        category_axis = inserted_chart.category_axis
        # category_axis.has_major_gridlines = True
        # category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
        category_axis.tick_labels.font.size = Pt(12)
        category_axis.tick_labels.font.name = "思源黑体 Light"
        value_axis = inserted_chart.value_axis
        value_axis.tick_labels.font.size = Pt(12)
        value_axis.tick_labels.font.name = "思源黑体 Light"
        chart_title = inserted_chart.chart_title
        chart_title.text_frame.fit_text(font_family="思源黑体 Normal", max_size=16)
        chart_title.text_frame.text = chart["title"]
        if chart["type"] == XL_CHART_TYPE.COLUMN_CLUSTERED:
            inserted_chart.series[0].format.fill.solid()
            inserted_chart.series[0].format.fill.fore_color.rgb = RGBColor(10, 182, 159)
        if chart["type"] == XL_CHART_TYPE.LINE:
            inserted_chart.series[0].format.line.fill.solid()
            inserted_chart.series[0].format.line.fill.fore_color.rgb = RGBColor(10, 182, 159)
        # chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        # chart_title.text_frame.paragraphs[0].font.name = "思源黑体 Normal"
        # chart_title.text_frame.text = 'test'

        # chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        # plot = inserted_chart.plots[0]
        # plot.has_data_labels = True
        # data_labels = plot.data_labels
        # data_labels.font.size = Pt(13)
        # data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        # data_labels.position = XL_LABEL_POSITION.INSIDE_END

    def generate_charts(self, df_report_item_data, language):
        # df = pandas.read_excel(r"C:\Users\zewen.liang\OneDrive - JCDECAUX\桌面\dummy data for imp rep.xlsx", parse_dates=['date'])
        df = df_report_item_data[["date", "impression"]]
        idx = pandas.date_range(df['date'].min(), df['date'].max())
        df = df.groupby('date', as_index=True)['impression'].sum()
        df = df.reindex(idx, fill_value=0)
        df = pandas.DataFrame(df)
        df = df.rename_axis('date')
        df.reset_index('date', inplace=True)

        df['year_month'] = df['date'].dt.strftime('%Y-%m')
        monthly_df = df.groupby('year_month', as_index=False)['impression'].sum()
        df['year_week'] = df['date'].dt.strftime('%Y-%U')
        weekly_df = df.groupby('year_week', as_index=False)['impression'].sum()
        daily_df = df.groupby('date', as_index=False)['impression'].sum()

        charts = {}

        # generating monthly bar chart
        monthly_bar_chart = {}
        chart_data = CategoryChartData()
        chart_data.categories = monthly_df["year_month"].tolist()
        chart_data.add_series('曝光量', monthly_df["impression"].tolist())
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        monthly_bar_chart.update({"data": chart_data})
        monthly_bar_chart.update({"type": chart_type})
        if language in ['cn', 'CN']:
            monthly_bar_chart.update({"title": "各年-月曝光量"})
        else:
            monthly_bar_chart.update({"title": "Each Year-Month Impression"})
        monthly_bar_chart.update({"object": "chart"})
        charts.update({"monthly bar chart": monthly_bar_chart})
        # generating weekly bar chart
        weekly_bar_chart = {}
        chart_data = CategoryChartData()
        chart_data.categories = weekly_df["year_week"].tolist()
        chart_data.add_series('曝光量', weekly_df["impression"].tolist())
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        weekly_bar_chart.update({"data": chart_data})
        weekly_bar_chart.update({"type": chart_type})
        if language in ['cn', 'CN']:
            weekly_bar_chart.update({"title": "各年-周曝光量"})
        else:
            weekly_bar_chart.update({"title": "Each Year-week Impression"})
        weekly_bar_chart.update({"object": "chart"})
        charts.update({"weekly bar chart": weekly_bar_chart})
        # generating daily line chart
        daily_line_chart = {}
        chart_data = ChartData()
        chart_data.categories = daily_df["date"].tolist()
        chart_data.add_series('曝光量', daily_df["impression"].tolist())
        chart_type = XL_CHART_TYPE.LINE
        daily_line_chart.update({"data": chart_data})
        daily_line_chart.update({"type": chart_type})
        if language in ['cn', 'CN']:
            daily_line_chart.update({"title": "按日曝光量趋势"})
        else:
            daily_line_chart.update({"title": "Daily Impression Trend"})
        daily_line_chart.update({"object": "chart"})
        charts.update({"daily line chart": daily_line_chart})

        return charts

    def generate_calendar_tables(self, df_report_item_data):
        df = df_report_item_data[["date", "impression"]]
        idx = pandas.date_range(df['date'].min(), df['date'].max())
        df = df.groupby('date', as_index=True)['impression'].sum()
        df = df.reindex(idx, fill_value=0)
        df = pandas.DataFrame(df)
        df = df.rename_axis('date')
        df.reset_index('date', inplace=True)

        df['year_month'] = df['date'].dt.strftime('%Y-%m')
        monthly_df = df.groupby('year_month', as_index=False)['impression'].sum()
        df['year_week'] = df['date'].dt.strftime('%Y-%U')
        weekly_df = df.groupby('year_week', as_index=False)['impression'].sum()
        daily_df = df.groupby('date', as_index=False)['impression'].sum()
        dataframe = daily_df.copy()
        ###
        dataframe['year_month'] = dataframe['date'].dt.strftime('%Y-%m')
        dataframe['day'] = dataframe['date'].dt.strftime('%d').astype(int)
        # print(dataframe['year_month'].unique())
        calendar_dicts = {}
        for year_month in dataframe['year_month'].unique():
            year, month = map(int, year_month.split("-"))
            date_obj = datetime(year, month, 1)
            total_days = (date_obj.replace(day=1, month=month % 12 + 1, year=(date_obj.year if date_obj.month != 12 else date_obj.year+1)) - date_obj).days
            # week_day = date_obj.weekday() # Monday=0, Sunday=6
            week_day = (date_obj.weekday() + 1) % 7 # Monday=1, Sunday=0
            # print(year_month, total_days, week_day)
            first_row_occupied = 7 - week_day
            rest_days = total_days - first_row_occupied
            rest_days_row_count = rest_days // 7 if rest_days % 7 == 0 else rest_days // 7 + 1
            total_row_count = rest_days_row_count + 1 + 1 # add a week day row (like Sunday, Monday etc)
            sub_df = dataframe[dataframe['year_month'] == year_month]

            calendar_cells = {1: {}}
            # for i, item in enumerate(['日','一','二','三','四','五','六']):
            #     calendar_cells += [[0, i, item]]
            # print(calendar_cells)
            row_no = 1
            column_no = week_day
            sub_dict = sub_df.set_index('day')['impression'].to_dict()
            for i in range(1, total_days + 1):
                impression = sub_dict[i].quantize(Decimal('0'), rounding=ROUND_HALF_UP) if i in sub_dict.keys() else 0
                calendar_cells[row_no].update({column_no: {"day": i, "impression": impression}})
                if column_no == 6:
                    row_no += 1
                    column_no = 0
                    calendar_cells.update({row_no: {}})
                else:
                    column_no += 1

            calendar_dicts.update({year_month: calendar_cells})
            tables = {}
            for item in calendar_dicts.keys():
                tables.update({item: {'object': 'table', 'table_title': item, 'table_dict': calendar_dicts[item]}})
        return tables

    def insert_table_into_content_placeholder(self, placeholder, table, shapes):
        # total_dict = generate_calendar_dicts(None)
        # table_dict = {'2023-10': total_dict['2023-10']}
        # table_dict = total_dict['2023-10']
        table_dict = table["table_dict"]
        table_title = table['table_title']
        idx = placeholder.placeholder_format.idx
        left = placeholder.left.cm
        top = placeholder.top.cm
        width = placeholder.width.cm
        height = placeholder.height.cm
        sp = placeholder.element
        sp.getparent().remove(sp)
        x, y, cx, cy = Cm(left), Cm(top), Cm(width), Cm(height)
        inserted_table = shapes.add_table(len(table_dict)+2, 7, x, y, cx, cy).table
        cell = inserted_table.cell(0, 0)
        cell.merge(inserted_table.cell(0, 6))
        cell.text = table_title
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].font.name = "思源黑体 Normal"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(244, 254, 253)
        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 设置垂直居中
        for i, item in enumerate(['日', '一', '二', '三', '四', '五', '六']):
            cell = inserted_table.cell(1, i)
            cell.text = item
            # cell.text_frame.fit_text(font_family="思源黑体 Normal", max_size=14)
            cell.text_frame.paragraphs[0].font.name = "思源黑体 Normal"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 设置垂直居中
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(8,146,126)
        for i in range(2, len(table_dict)+2):
            for j in range(7):
                cell = inserted_table.cell(i, j)
                # set background to transparent
                cell.fill.background()
                # cell.fill.solid()
                # cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # cell.merge(inserted_table.cell(0, 6))
        # cell.text = 'title'
        for row in table_dict.keys():
            for col in table_dict[row].keys():
                cell = inserted_table.cell(row+1, col)
                day = str(table_dict[row][col]['day'])
                impression = str(table_dict[row][col]['impression']) if table_dict[row][col]['impression'] !=0 else ""
                cell.text = day + "\n" + format(int(impression),',')
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(244, 254, 253)
                cell.text_frame.paragraphs[0].font.name = "思源黑体 Normal"
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(8, 146, 126)
                cell.text_frame.paragraphs[1].font.name = "思源黑体 Medium"
                cell.text_frame.paragraphs[1].font.size = Pt(14) if int(impression)<=99999 else Pt(12) if int(impression)<=9999999 else Pt(10)
                cell.text_frame.paragraphs[1].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 设置垂直居中

    def one_key_generate_report(self, report_id, language, env_info_code):
        if language in ['cn', 'CN']:
            prs = Presentation(cn_template_dir)
        else:
            prs = Presentation(en_template_dir)

        master = prs.slide_masters[0]
        layouts = master.slide_layouts
        layout_index = dict(zip([x.name for x in layouts], range(len(layouts))))

        # get dict_report_info_data and dict_env_company_info_data and df_report_item_data
        dict_report_info_data, dict_env_company_info_data, df_report_item_data =self.get_report_info(report_id, env_info_code)
        print(dict_report_info_data)
        # generate charts
        charts = self.generate_charts(df_report_item_data, language)
        #fill in 1st page
        if language in ['cn', 'CN']:
            report_main_title = "广告活动效果报告"
            report_subtitle = dict_report_info_data['campaign_name_cn']
            client_company_logo_dir = dict_report_info_data['client_logo_cn']
        else:
            report_main_title = "Campaign Performance Report"
            report_subtitle = dict_report_info_data['campaign_name_en']
            client_company_logo_dir = dict_report_info_data['client_logo_en']
        campaign_period = f"{dict_report_info_data['campaign_start_date']}~{dict_report_info_data['campaign_end_date']}"
        report_no = dict_report_info_data['no']
        company_logo_dir = dict_env_company_info_data['logo']
        prs = self.add_cover(prs, layouts[layout_index['Cover Shared']], report_main_title, report_subtitle, campaign_period,
                  report_no, company_logo_dir, client_company_logo_dir)
        # fill in 2nd page
        if language in ['cn', 'CN']:
            title_of_this_page = f"{dict_report_info_data['campaign_name_cn']}广告投放主要指标"
            descriptive_text = "本次投放主要指标如下："
        else:
            title_of_this_page = f"{dict_report_info_data['campaign_name_en']} Advertising Main Indicators"
            descriptive_text = "The main indicators of this campaign are as follows:"
        start_date = str(dict_report_info_data['campaign_start_date'])
        end_date = str(dict_report_info_data['campaign_end_date'])
        day_count = str((dict_report_info_data['campaign_end_date'] - dict_report_info_data['campaign_start_date']).days + 1)
        line_count = str(df_report_item_data['line_code'].nunique())
        station_count = str(df_report_item_data['station_code'].nunique())
        asset_count = str(df_report_item_data['asset_code'].nunique())
        impression = str(dict_report_info_data['total_impression'].quantize(Decimal('0'), rounding=ROUND_HALF_UP))
        cpm = str(dict_report_info_data['cpm'].quantize(Decimal('0.00'), rounding=ROUND_HALF_UP))
        prs = self.add_info_page(prs, layouts[layout_index["Info Page"]], title_of_this_page, descriptive_text, start_date, end_date, day_count, line_count, station_count,
                      asset_count, impression, cpm, company_logo_dir, client_company_logo_dir, report_no)
        # fill in 3rd page
        if language in ['cn', 'CN']:
            title_of_this_page = "曝光量按月及按周数据柱状图"
            descriptive_text = "本次活动曝光量按月分布图及按周分布图如下所示，您可以通过对比不同柱状图的高度了解相关月份或周的曝光效果："
        else:
            title_of_this_page = "Impression by month and by week data bar chart"
            descriptive_text = "The monthly and weekly distribution charts for the impression of this campaign are shown below. You can understand the impression effects of relevant months or weeks by comparing the heights of different bars:"
        prs = self.add_chart_page(prs, layouts[layout_index["Chart Page"]], title_of_this_page, descriptive_text,charts["monthly bar chart"], charts["weekly bar chart"], company_logo_dir, client_company_logo_dir, report_no)
        # generate calandar tables
        tables = self.generate_calendar_tables(df_report_item_data)

        # fill in 4th page
        if language in ['cn', 'CN']:
            title_of_this_page = "曝光量按日数据折线图及每日明细曝光量日历图"
            descriptive_text = "本次活动曝光量按日数据折线图及每日明细曝光量日历图如下所示，您可以查看每日明细曝光数值："
        else:
            title_of_this_page = "Daily impression line chart and calendar chart"
            descriptive_text = "The daily impression line chart and calendar chart are shown below, you can view the detailed impression value of each day:"
        prs = self.add_chart_page(prs, layouts[layout_index["Chart Page"]], title_of_this_page, descriptive_text, charts["daily line chart"], tables[list(tables.keys())[0]],
                       company_logo_dir, client_company_logo_dir, report_no)
        del tables[list(tables.keys())[0]]
        # fill in the rest pages of calendar charts
        if language in ['cn', 'CN']:
            title_of_this_page = "每日明细曝光量日历图"
            descriptive_text = "（续）"
        else:
            title_of_this_page = "Daily detailed impression calendar chart"
            descriptive_text = "(continued)"
        rest_placeholder_count = 0
        while len(tables.keys()) != 0:
            if rest_placeholder_count < len(tables.keys()):
                if len(tables.keys()) >= 2:
                    prs = self.add_chart_page(prs, layouts[layout_index["Chart Page"]], title_of_this_page, descriptive_text, tables[list(tables.keys())[0]],
                                   tables[list(tables.keys())[1]],
                                   company_logo_dir, client_company_logo_dir, report_no)
                    del tables[list(tables.keys())[0]]
                    del tables[list(tables.keys())[0]]

                if len(tables.keys()) == 1:
                    prs = self.add_chart_page(prs, layouts[layout_index["Chart Page"]], title_of_this_page, descriptive_text, tables[list(tables.keys())[0]], None,
                                   company_logo_dir, client_company_logo_dir, report_no)
                    del tables[list(tables.keys())[0]]
        # fill in confidentiality agreement
        # add_confidentiality_agreement(company_logo_dir, client_company_logo_dir, report_no)
        # fill in last page
        prs = self.add_last_page(prs, layouts[layout_index["Last Page"]], company_logo_dir, client_company_logo_dir, report_no)
        # delete layouts that have not been used to save storage
        for layout in layouts:
            try:
                layouts.remove(layout)
            except:
                pass
        # save report
        if language in ['cn', 'CN']:
            report_dir = rf".\{dict_report_info_data['client_name_cn']} {dict_report_info_data['no']}.pptx"
        else:
            report_dir = rf".\{dict_report_info_data['client_name_en']} {dict_report_info_data['no']}.pptx"
        prs.save(report_dir)
        return report_dir

# if __name__ == "__main__":
#     gr = generateReport()
#     report_dir = gr.one_key_generate_report(1, 'cn', 'mam-shanghai')
#     print(report_dir)

 
